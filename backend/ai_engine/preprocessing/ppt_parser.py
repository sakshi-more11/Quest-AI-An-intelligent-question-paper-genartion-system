"""
ppt_parser.py

Extract text from Microsoft PowerPoint (.pptx) presentations.
"""

from pathlib import Path
from pptx import Presentation


class PPTParser:

    def extract(self, ppt_path: str) -> dict:

        ppt_path = Path(ppt_path)

        if not ppt_path.exists():
            raise FileNotFoundError(f"{ppt_path} not found.")

        presentation = Presentation(ppt_path)

        slides = []

        for slide_number, slide in enumerate(presentation.slides, start=1):

            slide_text = []

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = shape.text.strip()

                    if text:
                        slide_text.append(text)

            slides.append({
                "page": slide_number,
                "text": "\n".join(slide_text)
            })

        full_text = "\n".join(
            slide["text"] for slide in slides
        )

        return {
            "filename": ppt_path.name,
            "file_type": "pptx",
            "pages": len(slides),
            "text": full_text,
            "page_data": slides
        }