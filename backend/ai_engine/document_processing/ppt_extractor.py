from pptx import Presentation


def extract_ppt(file_path: str):

    presentation = Presentation(file_path)

    text = []

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                if shape.text.strip():

                    text.append(shape.text)

    return "\n".join(text)